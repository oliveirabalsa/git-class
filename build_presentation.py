from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT / ".vendor"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BG = "F8F5EF"
CARD = "FFFDFC"
INK = "1F2A33"
MUTED = "66727C"
ACCENT = "2F7F79"
ACCENT_DARK = "205C57"
ACCENT_SOFT = "DCEEEB"
HIGHLIGHT = "F2A65A"
HIGHLIGHT_SOFT = "FBE8D2"
BLUE_SOFT = "DEEAF6"
RED_SOFT = "F8DDDA"
BORDER = "D9DFDA"
TERMINAL = "162029"
TERM_TEXT = "D8F6E7"


def rgb(hex_value):
    return RGBColor.from_string(hex_value)


def set_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    font_size=20,
    color=INK,
    bold=False,
    font_name="Arial",
    align=PP_ALIGN.LEFT,
    italic=False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return box


def add_bullets(
    slide,
    items,
    x,
    y,
    w,
    h,
    *,
    font_size=18,
    color=INK,
    bullet_color=ACCENT,
    font_name="Arial",
    space_after=4,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    frame.vertical_anchor = MSO_ANCHOR.TOP

    for idx, item in enumerate(items):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(space_after)
        bullet_run = paragraph.add_run()
        bullet_run.text = "• "
        bullet_run.font.name = font_name
        bullet_run.font.size = Pt(font_size)
        bullet_run.font.bold = True
        bullet_run.font.color.rgb = rgb(bullet_color)
        text_run = paragraph.add_run()
        text_run.text = item
        text_run.font.name = font_name
        text_run.font.size = Pt(font_size)
        text_run.font.color.rgb = rgb(color)
    return box


def add_card(slide, x, y, w, h, *, fill=CARD, line=BORDER, radius=True):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    return shape


def add_card_with_copy(
    slide,
    title,
    body,
    x,
    y,
    w,
    h,
    *,
    fill=CARD,
    title_color=INK,
    body_color=MUTED,
    accent=None,
    title_size=18,
    body_size=13,
):
    add_card(slide, x, y, w, h, fill=fill)
    if accent:
        strip = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.12)
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = rgb(accent)
        strip.line.fill.background()
    add_text(
        slide,
        title,
        x + 0.22,
        y + 0.2,
        w - 0.44,
        0.45,
        font_size=title_size,
        bold=True,
        color=title_color,
    )
    add_text(
        slide,
        body,
        x + 0.22,
        y + 0.7,
        w - 0.44,
        h - 0.9,
        font_size=body_size,
        color=body_color,
    )


def add_chip(slide, text, x, y, w, *, fill=ACCENT, color="FFFFFF"):
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34)
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = rgb(fill)
    chip.line.fill.background()
    add_text(
        slide,
        text,
        x + 0.08,
        y + 0.06,
        w - 0.16,
        0.2,
        font_size=11,
        color=color,
        bold=True,
    )


def add_terminal(slide, lines, x, y, w, h, *, title="Terminal"):
    block = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    block.fill.solid()
    block.fill.fore_color.rgb = rgb(TERMINAL)
    block.line.color.rgb = rgb("24313B")
    block.line.width = Pt(1)

    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = rgb(ACCENT_DARK)
    header.line.fill.background()
    add_text(
        slide,
        title,
        x + 0.16,
        y + 0.06,
        w - 0.32,
        0.2,
        font_size=12,
        color="FFFFFF",
        bold=True,
    )

    text = "\n".join(lines)
    box = slide.shapes.add_textbox(
        Inches(x + 0.22), Inches(y + 0.52), Inches(w - 0.44), Inches(h - 0.72)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Courier New"
    run.font.size = Pt(17)
    run.font.color.rgb = rgb(TERM_TEXT)
    return block


def add_arrow_text(slide, x, y, w, h, *, color=ACCENT, font_size=24):
    add_text(
        slide,
        "→",
        x,
        y,
        w,
        h,
        font_size=font_size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_footer(slide, page_num, total):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, Inches(0.55), Inches(7.02), Inches(12.75), Inches(7.02)
    )
    line.line.color.rgb = rgb(BORDER)
    line.line.width = Pt(1)
    add_text(slide, "Git do zero", 0.58, 7.08, 2.2, 0.18, font_size=10, color=MUTED)
    add_text(
        slide,
        f"{page_num}/{total}",
        12.0,
        7.08,
        0.6,
        0.18,
        font_size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def start_content_slide(prs, page_num, total, section, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_chip(slide, section.upper(), 0.55, 0.32, 1.6)
    add_text(slide, title, 0.55, 0.72, 8.8, 0.48, font_size=26, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.55, 1.18, 11.8, 0.28, font_size=12, color=MUTED)
    add_footer(slide, page_num, total)
    return slide


def build_cover(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, color="F4F0E7")

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.24)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(ACCENT)
    band.line.fill.background()

    add_chip(slide, "PT-BR | NIVEL INICIANTE", 0.7, 0.55, 2.25, fill=ACCENT_DARK)
    add_text(slide, "Git do zero", 0.7, 1.15, 4.6, 0.7, font_size=30, bold=True)
    add_text(
        slide,
        "Uma apresentacao simples para ensinar Git a quem nunca versionou nada.",
        0.72,
        1.95,
        4.9,
        0.52,
        font_size=16,
        color=MUTED,
    )
    add_bullets(
        slide,
        [
            "Explica conceitos sem jargao tecnico pesado.",
            "Mostra os primeiros comandos que realmente importam.",
            "Fecha com um exercicio curto para praticar em aula.",
        ],
        0.78,
        2.85,
        4.6,
        2.0,
        font_size=18,
    )

    add_card(slide, 6.0, 0.95, 6.15, 5.6, fill=CARD)
    add_text(slide, "Como o Git pensa", 6.35, 1.28, 2.6, 0.36, font_size=21, bold=True)
    add_text(
        slide,
        "Visual para repetir varias vezes durante a aula.",
        6.35,
        1.68,
        4.8,
        0.22,
        font_size=12,
        color=MUTED,
    )

    step_y = 2.15
    step_w = 1.55
    add_card_with_copy(
        slide,
        "1. Editar",
        "Voce mexe nos arquivos normalmente.",
        6.35,
        step_y,
        step_w,
        1.65,
        fill=ACCENT_SOFT,
        accent=ACCENT,
        body_size=12,
    )
    add_arrow_text(slide, 7.98, 2.55, 0.35, 0.45)
    add_card_with_copy(
        slide,
        "2. Preparar",
        "O comando git add separa o que vai entrar no proximo salvamento.",
        8.32,
        step_y,
        1.7,
        1.65,
        fill=HIGHLIGHT_SOFT,
        accent=HIGHLIGHT,
        body_size=11,
    )
    add_arrow_text(slide, 10.1, 2.55, 0.35, 0.45)
    add_card_with_copy(
        slide,
        "3. Commitar",
        "O commit cria um ponto do historico com mensagem.",
        10.45,
        step_y,
        1.35,
        1.65,
        fill=BLUE_SOFT,
        accent=ACCENT_DARK,
        body_size=11,
    )

    add_terminal(
        slide,
        ["git status", "git add README.md", 'git commit -m "Cria README"'],
        6.35,
        4.25,
        5.45,
        1.7,
        title="Comandos que aparecem na aula",
    )
    add_footer(slide, page_num, total)
    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Git do zero"
    prs.core_properties.subject = "Apresentacao introdutoria de Git em PT-BR"
    prs.core_properties.author = "OpenAI Codex"

    total = 13

    build_cover(prs, 1, total)

    slide = start_content_slide(
        prs,
        2,
        total,
        "Roteiro",
        "O que vamos aprender hoje",
        "Objetivo: sair do zero e entender o fluxo basico do Git.",
    )
    agenda_items = [
        ("O que e Git", "Por que ele existe e por que ajuda voce."),
        ("Repositorio", "O que significa iniciar um projeto com Git."),
        ("Commit", "Como salvar um ponto do historico."),
        ("Status e log", "Como ler o que esta acontecendo."),
        ("Branches", "Como testar sem baguncar a linha principal."),
        ("GitHub", "Como enviar e receber mudancas pela internet."),
    ]
    positions = [
        (0.7, 1.85),
        (4.45, 1.85),
        (8.2, 1.85),
        (0.7, 4.15),
        (4.45, 4.15),
        (8.2, 4.15),
    ]
    fills = [CARD, HIGHLIGHT_SOFT, BLUE_SOFT, ACCENT_SOFT, CARD, HIGHLIGHT_SOFT]
    for (title, body), (x, y), fill in zip(agenda_items, positions, fills):
        add_card_with_copy(slide, title, body, x, y, 3.45, 1.75, fill=fill, accent=ACCENT)

    slide = start_content_slide(
        prs,
        3,
        total,
        "Conceito",
        "O que e Git?",
        "Explique com uma analogia simples antes de mostrar qualquer comando.",
    )
    add_card(slide, 0.72, 1.75, 6.0, 4.6, fill=CARD)
    add_text(slide, "Git guarda versoes do seu projeto.", 1.0, 2.05, 5.2, 0.4, font_size=23, bold=True)
    add_text(
        slide,
        "Pense nele como um historico inteligente: voce salva momentos importantes, consegue voltar atras e enxerga exatamente o que mudou.",
        1.0,
        2.6,
        5.2,
        1.2,
        font_size=17,
        color=MUTED,
    )
    add_card_with_copy(
        slide,
        "Frase boa para a aula",
        "Git nao e uma pasta magica. Ele e um historico organizado da sua pasta.",
        1.0,
        4.2,
        5.1,
        1.55,
        fill=ACCENT_SOFT,
        accent=ACCENT,
        body_size=15,
    )
    right_cards = [
        ("Historico", "Voce sabe quem mudou, quando mudou e o que mudou.", HIGHLIGHT_SOFT),
        ("Seguranca", "Se der ruim, fica mais facil voltar para um estado bom.", BLUE_SOFT),
        ("Colaboracao", "Mais de uma pessoa consegue trabalhar sem caos.", CARD),
    ]
    for idx, (title, body, fill) in enumerate(right_cards):
        add_card_with_copy(
            slide,
            title,
            body,
            7.15,
            1.9 + (idx * 1.48),
            5.3,
            1.18,
            fill=fill,
            accent=ACCENT,
            body_size=12,
        )

    slide = start_content_slide(
        prs,
        4,
        total,
        "Base",
        "Tres conceitos que precisam ficar claros",
        "Use este slide como mapa mental do resto da apresentacao.",
    )
    card_data = [
        ("Working directory", "Onde voce edita os arquivos agora.", ACCENT_SOFT),
        ("Stage", "Area de preparo do proximo commit.", HIGHLIGHT_SOFT),
        ("Commit", "Foto com legenda do projeto naquele momento.", BLUE_SOFT),
    ]
    x_positions = [0.8, 4.45, 8.1]
    for (title, body, fill), x in zip(card_data, x_positions):
        add_card_with_copy(slide, title, body, x, 2.3, 3.05, 2.15, fill=fill, accent=ACCENT, title_size=20, body_size=15)
    add_arrow_text(slide, 3.95, 3.0, 0.3, 0.5)
    add_arrow_text(slide, 7.6, 3.0, 0.3, 0.5)
    add_card_with_copy(
        slide,
        "Regra mental",
        "Editar → preparar → salvar. Se o aluno decorar isso, metade da aula ja funcionou.",
        1.65,
        5.2,
        10.0,
        0.95,
        fill=CARD,
        accent=ACCENT_DARK,
        body_size=14,
    )

    slide = start_content_slide(
        prs,
        5,
        total,
        "Inicio",
        "Antes de comecar no terminal",
        "Mostre so o necessario para a pessoa se localizar.",
    )
    add_card_with_copy(
        slide,
        "Checklist rapido",
        "1. Instale o Git.\n2. Abra o terminal.\n3. Entre na pasta da aula.\n4. Confirme que o Git responde.",
        0.8,
        1.9,
        4.25,
        3.2,
        fill=CARD,
        accent=ACCENT,
        body_size=16,
    )
    add_bullets(
        slide,
        [
            "Use git --version para ver se a instalacao esta ok.",
            "Use pwd para saber em que pasta voce esta.",
            "Use ls para listar os arquivos da pasta atual.",
        ],
        0.98,
        5.25,
        4.0,
        1.2,
        font_size=16,
    )
    add_terminal(slide, ["git --version", "pwd", "ls"], 5.45, 1.95, 6.8, 2.55, title="Primeiros comandos")
    add_card_with_copy(
        slide,
        "O que voce quer ouvir do computador",
        "Se aparecer um numero de versao, o Git esta pronto para a aula.",
        5.45,
        4.8,
        6.8,
        1.2,
        fill=ACCENT_SOFT,
        accent=ACCENT_DARK,
        body_size=15,
    )

    slide = start_content_slide(
        prs,
        6,
        total,
        "Pratica",
        "Criando seu primeiro repositorio",
        "Aqui o aluno ve pela primeira vez o git init em contexto real.",
    )
    step_titles = [
        "Criar pasta",
        "Entrar na pasta",
        "Iniciar Git",
        "Conferir status",
    ]
    step_commands = ["mkdir aula-git", "cd aula-git", "git init", "git status"]
    for idx, (title, command) in enumerate(zip(step_titles, step_commands)):
        y = 1.8 + idx * 1.1
        add_card_with_copy(
            slide,
            title,
            command,
            0.82,
            y,
            4.15,
            0.88,
            fill=CARD if idx % 2 == 0 else HIGHLIGHT_SOFT,
            accent=ACCENT,
            title_size=17,
            body_size=15,
        )
    add_terminal(
        slide,
        ["mkdir aula-git", "cd aula-git", "git init", "git status"],
        5.45,
        1.8,
        6.75,
        3.3,
        title="Sequencia completa",
    )
    add_card_with_copy(
        slide,
        "Leitura esperada",
        "Depois do git init, o Git cria um repositorio vazio. Depois do git status, ele diz que ainda nao existe nada para commitar.",
        5.45,
        5.25,
        6.75,
        1.05,
        fill=BLUE_SOFT,
        accent=ACCENT_DARK,
        body_size=13,
    )

    slide = start_content_slide(
        prs,
        7,
        total,
        "Pratica",
        "Como salvar mudancas",
        "Este e o ciclo principal que voce vai repetir o tempo todo.",
    )
    flow_data = [
        ("1. Editar", "Voce cria ou altera um arquivo.", ACCENT_SOFT),
        ("2. git add", "Voce escolhe o que vai entrar no proximo salvamento.", HIGHLIGHT_SOFT),
        ("3. git commit", "Voce grava um ponto do historico com mensagem.", BLUE_SOFT),
    ]
    flow_x = [0.85, 4.45, 8.05]
    for (title, body, fill), x in zip(flow_data, flow_x):
        add_card_with_copy(slide, title, body, x, 1.88, 3.0, 1.6, fill=fill, accent=ACCENT, title_size=19, body_size=13)
    add_arrow_text(slide, 3.96, 2.4, 0.28, 0.4)
    add_arrow_text(slide, 7.57, 2.4, 0.28, 0.4)
    add_terminal(
        slide,
        [
            'echo "# Aula Git" > README.md',
            "git add README.md",
            'git commit -m "Cria README inicial"',
        ],
        0.9,
        4.15,
        6.1,
        1.95,
        title="Exemplo para mostrar ao vivo",
    )
    add_card_with_copy(
        slide,
        "Mensagem boa de commit",
        "Escreva o que mudou de forma curta e direta. Exemplo: Cria README inicial.",
        7.35,
        4.15,
        4.75,
        1.95,
        fill=CARD,
        accent=HIGHLIGHT,
        body_size=15,
    )

    slide = start_content_slide(
        prs,
        8,
        total,
        "Leitura",
        "Os 3 comandos que voce mais vai usar",
        "Status, log e diff ajudam o aluno a ganhar confianca.",
    )
    command_cards = [
        ("git status", "Mostra o que mudou e o que ainda falta salvar.", ACCENT_SOFT),
        ("git log --oneline", "Mostra o historico em formato curto e facil de ler.", BLUE_SOFT),
        ("git diff", "Mostra linha por linha o que mudou no arquivo.", HIGHLIGHT_SOFT),
    ]
    for idx, (title, body, fill) in enumerate(command_cards):
        add_card_with_copy(slide, title, body, 0.85 + idx * 4.1, 2.0, 3.5, 2.2, fill=fill, accent=ACCENT, title_size=18, body_size=14)
    add_terminal(
        slide,
        ["git status", "git diff", "git log --oneline"],
        2.35,
        4.7,
        8.55,
        1.25,
        title="Comandos para repetir varias vezes",
    )

    slide = start_content_slide(
        prs,
        9,
        total,
        "Branches",
        "Branches sem drama",
        "Apresente branch como um caminho paralelo, nao como um misterio.",
    )
    add_card_with_copy(
        slide,
        "Analogia que funciona",
        "Uma branch e um rascunho paralelo. Voce testa uma ideia sem mexer direto na linha principal do projeto.",
        0.85,
        1.95,
        5.0,
        2.1,
        fill=CARD,
        accent=ACCENT,
        title_size=21,
        body_size=15,
    )
    add_bullets(
        slide,
        [
            "Crie branch para experimentar com seguranca.",
            "Volte para main quando quiser comparar ou seguir o fluxo normal.",
            "No mundo real, depois voce junta tudo com merge ou pull request.",
        ],
        1.05,
        4.45,
        4.75,
        1.65,
        font_size=16,
    )
    add_terminal(
        slide,
        ["git switch -c nova-ideia", "git switch main"],
        6.35,
        2.05,
        5.65,
        1.85,
        title="Comandos minimos para ensinar",
    )
    add_card_with_copy(
        slide,
        "Ideia principal",
        "No nivel iniciante, basta entender que branch = caminho paralelo de trabalho.",
        6.35,
        4.35,
        5.65,
        1.55,
        fill=HIGHLIGHT_SOFT,
        accent=HIGHLIGHT,
        body_size=15,
    )

    slide = start_content_slide(
        prs,
        10,
        total,
        "GitHub",
        "Git e GitHub: como se conectam",
        "Mantenha a explicacao binaria: Git e local, GitHub e remoto.",
    )
    add_card_with_copy(
        slide,
        "Git",
        "Roda no seu computador. E onde voce faz add, commit, log e branch.",
        0.9,
        1.95,
        3.7,
        1.65,
        fill=ACCENT_SOFT,
        accent=ACCENT_DARK,
        title_size=24,
        body_size=14,
    )
    add_arrow_text(slide, 4.75, 2.45, 0.4, 0.5)
    add_card_with_copy(
        slide,
        "GitHub",
        "Guarda o repositorio na internet para compartilhar, enviar e receber mudancas.",
        5.2,
        1.95,
        4.8,
        1.65,
        fill=BLUE_SOFT,
        accent=ACCENT,
        title_size=24,
        body_size=14,
    )
    add_arrow_text(slide, 10.1, 2.45, 0.4, 0.5)
    add_card_with_copy(
        slide,
        "Traduza assim",
        "push envia. pull traz.",
        10.55,
        1.95,
        1.85,
        1.65,
        fill=HIGHLIGHT_SOFT,
        accent=HIGHLIGHT,
        title_size=18,
        body_size=15,
    )
    add_terminal(
        slide,
        [
            "git remote add origin URL_DO_REPOSITORIO",
            "git push -u origin main",
            "git pull",
        ],
        1.45,
        4.2,
        10.6,
        1.95,
        title="Comandos que valem decorar",
    )

    slide = start_content_slide(
        prs,
        11,
        total,
        "Rotina",
        "Fluxo simples do dia a dia",
        "Se o aluno decorar esta sequencia, ele ja consegue trabalhar em projetos pequenos.",
    )
    steps = ["git pull", "editar", "git status", "git add", "git commit", "git push"]
    fills = [BLUE_SOFT, CARD, ACCENT_SOFT, HIGHLIGHT_SOFT, CARD, BLUE_SOFT]
    for idx, (step, fill) in enumerate(zip(steps, fills), start=1):
        x = 0.62 + (idx - 1) * 2.1
        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(2.15), Inches(0.48), Inches(0.48)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = rgb(ACCENT)
        circle.line.fill.background()
        add_text(slide, str(idx), x + 0.11, 2.23, 0.25, 0.18, font_size=12, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        add_card_with_copy(slide, step, "", x + 0.62, 2.0, 1.35, 0.82, fill=fill, accent=ACCENT, title_size=16, body_size=1)
        if idx < len(steps):
            add_arrow_text(slide, x + 1.97, 2.13, 0.2, 0.3, font_size=18)
    add_card_with_copy(
        slide,
        "Resumo do ciclo",
        "Primeiro atualize, depois trabalhe em blocos pequenos, salve com commit e so entao envie para o remoto.",
        1.1,
        4.1,
        11.0,
        1.45,
        fill=CARD,
        accent=ACCENT_DARK,
        body_size=16,
    )

    slide = start_content_slide(
        prs,
        12,
        total,
        "Cuidados",
        "Erros comuns de iniciantes",
        "Um slide curto para evitar tropeços nas primeiras semanas.",
    )
    warnings = [
        ("Nao olhar git status", "Antes de add e commit, leia o estado do projeto.", RED_SOFT),
        ('Mensagem "update"', "Prefira dizer o que mudou de verdade.", HIGHLIGHT_SOFT),
        ("Ficar horas sem commit", "Commits pequenos ajudam a voltar atras.", CARD),
        ("Confundir Git com GitHub", "Git e local. GitHub e remoto.", ACCENT_SOFT),
    ]
    positions = [(0.85, 1.95), (6.75, 1.95), (0.85, 4.1), (6.75, 4.1)]
    for (title, body, fill), (x, y) in zip(warnings, positions):
        add_card_with_copy(slide, title, body, x, y, 5.1, 1.6, fill=fill, accent=ACCENT, title_size=18, body_size=14)
    add_text(
        slide,
        "Dica final: commits pequenos e frequentes quase sempre sao melhores do que um commit gigante no fim do dia.",
        1.15,
        6.2,
        11.0,
        0.28,
        font_size=13,
        color=MUTED,
        italic=True,
    )

    slide = start_content_slide(
        prs,
        13,
        total,
        "Fechamento",
        "Mini exercicio para praticar hoje",
        "Feche a aula com uma tarefa curta para consolidar o fluxo.",
    )
    add_card_with_copy(
        slide,
        "Tarefa",
        "1. Crie a pasta pratica-git.\n2. Rode git init.\n3. Crie um README.md.\n4. Rode git add e git commit.\n5. Altere o README e use git status, git diff e git log --oneline.",
        0.82,
        1.95,
        6.1,
        3.75,
        fill=CARD,
        accent=ACCENT,
        title_size=22,
        body_size=16,
    )
    add_card_with_copy(
        slide,
        "Ao final, a pessoa deve conseguir responder",
        "• O que esta staged agora?\n• O que ja foi commitado?\n• Como voltar para a branch main?\n• Qual comando envia mudancas para o GitHub?",
        7.3,
        1.95,
        5.0,
        2.6,
        fill=ACCENT_SOFT,
        accent=ACCENT_DARK,
        title_size=18,
        body_size=15,
    )
    add_card_with_copy(
        slide,
        "Se o aluno entendeu status, add, commit e push, ele ja saiu do zero.",
        "",
        7.3,
        4.95,
        5.0,
        0.95,
        fill=HIGHLIGHT_SOFT,
        accent=HIGHLIGHT,
        title_size=16,
        body_size=1,
    )

    out_dir = ROOT / "output"
    out_dir.mkdir(exist_ok=True)
    output_file = out_dir / "git-para-iniciantes-ptbr.pptx"
    prs.save(output_file)
    return output_file


if __name__ == "__main__":
    generated = build_deck()
    print(generated)
